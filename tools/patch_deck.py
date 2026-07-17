"""Apply post-rebuild fixes to the already-badged deck, in place.

The deck already has its badges, link band, and notes (from edit_deck.py).
This only patches what changed when the webapp was rebuilt:
  1. Title link band: placeholder URL -> real GitHub Pages URL.
  2. Slide 3 note: Attention is now a preset-only schematic, not live-on-chat.
  3. Slide 8 note: match the actual "Confidently wrong" preset button label.

Safe to run once. It edits text in place and adds no shapes.
"""
from pptx import Presentation
from pptx.util import Pt

PATH = r"c:\Users\akash.sood\OneDrive - City of Dallas-5516124-G Tenant\Documents\AI at COD\Transformers_Staff\Transformers_DBI_Staff.pptx"

REAL_URL = "akashsd.github.io/transformers_dbi_staff"

NOTE_ATTENTION = (
    "Open the webapp - Attention panel. Click 'she' and watch the lines fan out; the "
    "thickest one lands on 'officer'. There are two preset sentences - try the ambiguous "
    "'it' one too. Want more? Have people drop a pronoun sentence in chat and we'll reason "
    "through where it should point."
)

NOTE_NEXT_TOKEN = (
    "Open the Next-token panel. Click the 'Confidently wrong' preset - the Dallas Cowboys "
    "one. Watch it offer recent years with total confidence, every one wrong. This is why we "
    "don't ship these unsupervised. Then: 1 in chat if you've hit this at work, 2 if not yet."
)


def patch_title_band(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        paras = shape.text_frame.paragraphs
        if paras and "OPEN THIS WHILE I TALK" in paras[0].text and len(paras) > 1:
            run = paras[1].runs[0]
            run.text = REAL_URL
            run.font.size = Pt(14)  # real URL is longer; shrink so it fits the 4.5" band
            return True
    return False


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def main():
    prs = Presentation(PATH)
    slides = list(prs.slides)
    ok = patch_title_band(slides[0])
    print("title band URL:", "updated" if ok else "NOT FOUND")
    set_notes(slides[2], NOTE_ATTENTION)   # slide 3 (index 2)
    set_notes(slides[7], NOTE_NEXT_TOKEN)  # slide 8 (index 7)
    print("notes updated: slide 3 (attention), slide 8 (next-token)")
    prs.save(PATH)
    print("saved")


if __name__ == "__main__":
    main()
