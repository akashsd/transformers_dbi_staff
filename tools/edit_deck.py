from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PATH = r"c:\Users\akash.sood\OneDrive - City of Dallas-5516124-G Tenant\Documents\AI at COD\Transformers_Staff\Transformers_DBI_Staff.pptx"

NAVY = RGBColor(0x00, 0x1B, 0x3B)
BLUE = RGBColor(0x4D, 0xA3, 0xFF)
GREEN = RGBColor(0x8E, 0xC9, 0x00)
ORANGE = RGBColor(0xFF, 0xA9, 0x40)
LIGHT = RGBColor(0xEB, 0xF6, 0xFF)

BADGES = {
    2: ("📱", "TRY IT", "Tokenizer panel", BLUE),
    3: ("📱", "TRY IT", "Attention viz", BLUE),
    4: ("📊", "ZOOM POLL", "Spot the transformer", ORANGE),
    5: ("📊", "ZOOM POLL", "Which flavor for chat?", ORANGE),
    6: ("💬", "CHAT DROP", "What would you try?", GREEN),
    7: ("📱", "TRY IT", "Semantic search - live", BLUE),
    8: ("📱", "TRY IT", "Next-token predictor", BLUE),
    9: ("💬", "CHAT DROP", "Notebooks link", GREEN),
}

NOTES = {
    1: "Before we start - chat waterfall. Type one word for what you hope to leave with today, but DON'T send yet. Three, two, one - send.",
    2: "Open the webapp - Tokenizer panel. Guess how many tokens 'strawberry' is. Put your guess in chat, then check. I'll wait.",
    3: "Drop a sentence in chat with a pronoun in it. I'll pick two and run them through the attention viz live.",
    4: "Zoom poll launching now - which of these is a transformer? Multi-select. Take 15 seconds.",
    5: "Zoom poll: for a chatbot, which flavor? Vote before I reveal.",
    6: "Chat drop: one thing you'd want to try at DBI with this. Don't overthink it. I'll read a few aloud.",
    7: "This is the main event. Everyone in the Semantic Search panel? Good. Watch me run this first, then I'll take queries from chat. Bonus points if your query doesn't share a single word with the docs you want to find.",
    8: "Open the Next-token panel. Click the preset that says 'Dallas Cowboys.' Watch the model be confidently wrong. This is why we don't ship these unsupervised. Then: 1 in chat if you've hit one of these problems at work, 2 if not yet.",
    9: "Reposting the repo link in chat now - three notebooks in there, first one takes ten minutes and gives you a working semantic search over your own CSV.",
}


def add_badge(slide, icon, top_line, bottom_line, fill):
    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(10.0),
        Inches(0.42),
        Inches(2.75),
        Inches(0.55),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = fill
    badge.line.color.rgb = fill
    badge.line.width = Pt(0.5)
    tf = badge.text_frame
    tf.clear()
    tf.word_wrap = False

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run()
    r1.text = f"{icon}  {top_line}"
    r1.font.size = Pt(10)
    r1.font.bold = True
    r1.font.color.rgb = NAVY
    try:
        r1.font.character_spacing = Pt(3)
    except Exception:
        pass

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = bottom_line
    r2.font.size = Pt(11)
    r2.font.color.rgb = NAVY


def add_title_band(slide):
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.8),
        Inches(6.3),
        Inches(4.5),
        Inches(0.9),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(0x0E, 0x2A, 0x4A)
    band.line.color.rgb = BLUE
    band.line.width = Pt(1)

    tf = band.text_frame
    tf.clear()

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run()
    r1.text = "📱  OPEN THIS WHILE I TALK"
    r1.font.size = Pt(11)
    r1.font.bold = True
    r1.font.color.rgb = BLUE
    try:
        r1.font.character_spacing = Pt(4)
    except Exception:
        pass

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = "dbi-transformers.pages.dev"
    r2.font.size = Pt(18)
    r2.font.bold = True
    r2.font.color.rgb = LIGHT


def set_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.clear()
    notes_tf.text = text


def main():
    prs = Presentation(PATH)
    for index, slide in enumerate(prs.slides, start=1):
        if index == 1:
            add_title_band(slide)
        if index in BADGES:
            icon, top_line, bottom_line, fill = BADGES[index]
            add_badge(slide, icon, top_line, bottom_line, fill)
        if index in NOTES:
            set_notes(slide, NOTES[index])
    prs.save(PATH)
    print(f"updated {len(prs.slides)} slides")


if __name__ == "__main__":
    main()